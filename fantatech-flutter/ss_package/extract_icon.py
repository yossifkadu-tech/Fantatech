# -*- coding: utf-8 -*-
import sys, os, zipfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from PIL import Image, ImageDraw, ImageFilter

pptx_path = r"C:\Users\My laptop\Downloads\FantaTech icon.pptx"
out_dir   = r"C:\Users\My laptop\Desktop\smarthome-hub\fantatech-flutter\assets\images"
prev_dir  = r"C:\Users\My laptop\Desktop\smarthome-hub\fantatech-flutter"

os.makedirs(out_dir, exist_ok=True)

# ── 1. Extract embedded images ────────────────────────────────────────────
print("Extracting media...")
extracted = []
with zipfile.ZipFile(pptx_path, 'r') as z:
    media_files = sorted([f for f in z.namelist() if f.startswith('ppt/media/')])
    print(f"  {len(media_files)} file(s) found")
    for i, mf in enumerate(media_files):
        info = z.getinfo(mf)
        ext  = os.path.splitext(mf)[1].lower()
        print(f"  [{i}] {os.path.basename(mf)}  ({info.file_size:,} bytes)")
        data = z.read(mf)
        save_path = os.path.join(out_dir, f"icon_extract_{i}{ext}")
        with open(save_path, 'wb') as fp:
            fp.write(data)
        extracted.append((save_path, info.file_size))
        print(f"      saved -> {save_path}")

# ── 2. Slide structure ────────────────────────────────────────────────────
prs = Presentation(pptx_path)
print(f"\nSlides: {len(prs.slides)}")
for si, slide in enumerate(prs.slides):
    print(f"  Slide {si+1}: {len(slide.shapes)} shapes")
    for sh in slide.shapes:
        print(f"    {sh.name}  type={sh.shape_type}  "
              f"({int(sh.left/914400*96)}px,{int(sh.top/914400*96)}px) "
              f"{int(sh.width/914400*96)}x{int(sh.height/914400*96)}px")

# ── 3. Build comparison preview ───────────────────────────────────────────
# Pick the largest extracted image as the new icon source
if extracted:
    best = max(extracted, key=lambda x: x[1])
    src_path = best[0]
    print(f"\nUsing {src_path} as icon source")

    src_img = Image.open(src_path).convert("RGBA")
    print(f"Source size: {src_img.size}")

    def make_icon(img, size, radius_pct=0.22):
        """Resize to square (letterbox/center-crop), add rounded corners."""
        s = size
        r = int(s * radius_pct)
        # center-crop to square
        w, h = img.size
        side = min(w, h)
        left = (w - side) // 2
        top  = (h - side) // 2
        crop = img.crop((left, top, left+side, top+side))
        icon = crop.resize((s, s), Image.LANCZOS)
        # rounded corners mask
        mask = Image.new("L", (s, s), 0)
        d    = ImageDraw.Draw(mask)
        d.rounded_rectangle([0, 0, s-1, s-1], radius=r, fill=255)
        result = Image.new("RGBA", (s, s), (0, 0, 0, 0))
        result.paste(icon, mask=mask)
        return result

    # Generate preview: current icon vs new icon at multiple sizes
    SIZES   = [1024, 512, 192, 128, 64, 48]
    LABELS  = ["1024", "512", "192", "128", "64", "48"]
    padding = 20
    cell_w  = 1024 + padding * 2
    row_h   = 100 + padding

    # Load current icon
    cur_path = os.path.join(out_dir, "app_icon.png")
    cur_img  = Image.open(cur_path).convert("RGBA") if os.path.exists(cur_path) else None

    # Build preview canvas
    # Row 1: all sizes of NEW icon
    # Row 2: all sizes of CURRENT icon (for comparison)
    total_w = sum(s + padding for s in SIZES) + padding
    total_h = (max(SIZES) + row_h) * 2 + padding * 3 + 60

    canvas = Image.new("RGB", (total_w, total_h), (18, 18, 30))
    draw   = ImageDraw.Draw(canvas)

    def paste_row(label, img_src, y_offset, bg=(18,18,30)):
        draw.text((padding, y_offset), label, fill=(180, 180, 200))
        x = padding
        for s in SIZES:
            icon = make_icon(img_src, s)
            # center vertically in the max-size row
            iy = y_offset + 24 + (max(SIZES) - s) // 2
            # paste on background swatch
            bg_tile = Image.new("RGB", (s, s), bg)
            bg_tile.paste(icon, mask=icon.split()[3])
            canvas.paste(bg_tile, (x, iy))
            draw.text((x + s//2 - 8, iy + s + 4), f"{s}px", fill=(120,120,140))
            x += s + padding

    paste_row("NEW ICON  (FantaTech icon.pptx)", src_img, padding)
    if cur_img:
        paste_row("CURRENT ICON  (app_icon.png)", cur_img,
                  max(SIZES) + row_h + padding * 2)

    prev_path = os.path.join(prev_dir, "icon_preview.png")
    canvas.save(prev_path)
    print(f"Preview saved -> {prev_path}")

    # Also save the new source for flutter_launcher_icons
    new_src = os.path.join(out_dir, "app_icon_new.png")
    src_img.save(new_src)
    print(f"New icon source -> {new_src}")
